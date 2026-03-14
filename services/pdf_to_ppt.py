import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches


def pdf_to_ppt(input_path, output_path, dpi=200):
    """Convert PDF pages to a PowerPoint presentation (each page becomes a slide image)."""
    images = convert_from_path(input_path, dpi=dpi)

    prs = Presentation()
    # Set slide dimensions to match first page aspect ratio
    if images:
        first = images[0]
        w, h = first.size
        aspect = w / h
        slide_h = Inches(7.5)
        slide_w = int(slide_h * aspect)
        prs.slide_width = slide_w
        prs.slide_height = slide_h

    blank_layout = prs.slide_layouts[6]  # blank layout
    temp_dir = os.path.dirname(output_path)

    for i, img in enumerate(images):
        slide = prs.slides.add_slide(blank_layout)
        img_path = os.path.join(temp_dir, f"_temp_slide_{i}.png")
        img.save(img_path, 'PNG')

        slide.shapes.add_picture(
            img_path,
            left=0, top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )
        os.remove(img_path)

    prs.save(output_path)
    return output_path, len(images)
